"""
VaR Calculation & Evolution — Group Presentation Builder
Market Risk Modelling, Frankfurt School of Finance & Management
Dr. Sebastian Irle, Summer Term 2026

Usage:
    pip install python-pptx matplotlib
    python create_presentation.py

Output: var_presentation.pptx

Design:
    - Light/clean: white background, charcoal (#2D2D2D) accents
    - Font: Aptos throughout
    - Left vertical accent bar carried across all slides (consistent motif)
    - Two-column layout for each VaR method slide (Theory | Results)
    - Mathematical formulas rendered as crisp PNG images via matplotlib mathtext
    - Grey placeholder boxes with descriptive labels for missing figures

Figure naming convention (put PNGs in ./figures/):
    delta_normal_backtest.png  histsim_backtest.png  mc_backtest.png
    garch_backtest.png         evt_backtest.png
    qq_plots.png               acf_squared.png       arch_pvalues.png
"""

import io
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─────────────────────────────────────────────────────────────────────────────
# CONSTANTS
# ─────────────────────────────────────────────────────────────────────────────

SW = Inches(10)       # slide width  (16:9)
SH = Inches(5.625)    # slide height

C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK     = RGBColor(0x2D, 0x2D, 0x2D)
C_BODY     = RGBColor(0x3D, 0x3D, 0x3D)
C_MUTED    = RGBColor(0x88, 0x88, 0x88)
C_RULE     = RGBColor(0xCC, 0xCC, 0xCC)
C_PH_BG    = RGBColor(0xE8, 0xE8, 0xE8)
C_PH_TXT   = RGBColor(0x99, 0x99, 0x99)
C_ROW_ALT  = RGBColor(0xF5, 0xF5, 0xF5)
C_CARD_BG  = RGBColor(0xF4, 0xF4, 0xF4)
C_FORMULA  = RGBColor(0xF6, 0xF6, 0xF6)

FONT    = "Aptos"
FIG_DIR = Path("figures")
BAR_W   = 0.18


# ─────────────────────────────────────────────────────────────────────────────
# FORMULA RENDERING  (matplotlib mathtext → transparent PNG → BytesIO)
# ─────────────────────────────────────────────────────────────────────────────

def _formula_image(latex: str, font_size: int = 18, dpi: int = 200,
                   color: str = "#2D2D2D", pad: float = 0.10):
    """
    Auto-size a LaTeX formula to a transparent PNG.

    Returns (BytesIO, width_inches, height_inches).
    """
    # Pass 1 — measure text extents
    probe = plt.figure(figsize=(8, 2), dpi=dpi)
    probe.patch.set_alpha(0)
    t = probe.text(0.5, 0.5, f"${latex}$",
                   fontsize=font_size, ha="center", va="center", color=color)
    probe.canvas.draw()
    bb = t.get_window_extent(probe.canvas.get_renderer())
    plt.close(probe)

    pad_px = int(dpi * pad)
    w_in   = (bb.width  + pad_px * 2) / dpi
    h_in   = (bb.height + pad_px * 2) / dpi

    # Pass 2 — render at exact size
    fig, ax = plt.subplots(figsize=(w_in, h_in), dpi=dpi)
    fig.patch.set_alpha(0)
    ax.set_axis_off()
    ax.text(0.5, 0.5, f"${latex}$",
            fontsize=font_size, ha="center", va="center",
            color=color, transform=ax.transAxes)

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi,
                bbox_inches="tight", transparent=True, pad_inches=pad)
    plt.close(fig)
    buf.seek(0)
    return buf, w_in, h_in


def add_formula(slide, latex: str,
                center_x: float, top: float,
                max_w: float = 4.1, font_size: int = 18, dpi: int = 200):
    """
    Render a formula and place it horizontally centred around center_x.
    Scales down the font if the rendered width exceeds max_w.
    """
    buf, w, h = _formula_image(latex, font_size=font_size, dpi=dpi)

    if w > max_w:
        # Re-render at a reduced font size so it fits
        scale     = max_w / w
        font_size = max(10, int(font_size * scale * 0.9))
        buf, w, h = _formula_image(latex, font_size=font_size, dpi=dpi)
        if w > max_w:
            w = max_w

    left = center_x - w / 2
    slide.shapes.add_picture(buf, Inches(left), Inches(top),
                             Inches(w), Inches(h))


# ─────────────────────────────────────────────────────────────────────────────
# LOW-LEVEL HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill, line_color=None):
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             size=13, bold=False, italic=False, color=C_BODY,
             align=PP_ALIGN.LEFT, wrap=True,
             v_anchor=MSO_ANCHOR.TOP):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_bullets(slide, items, left, top, width, height,
                size=12, color=C_BODY, spacing_after=5):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(spacing_after)
        run = p.add_run()
        run.text = f"\u2013  {item}"
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return txb


def add_placeholder(slide, left, top, width, height, label):
    add_rect(slide, left, top, width, height, C_PH_BG)
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name = FONT
    run.font.size = Pt(10.5)
    run.font.italic = True
    run.font.color.rgb = C_PH_TXT


def add_figure_or_ph(slide, fig_name, left, top, width, height, label):
    path = FIG_DIR / fig_name
    if path.exists():
        slide.shapes.add_picture(
            str(path), Inches(left), Inches(top), Inches(width), Inches(height)
        )
    else:
        add_placeholder(slide, left, top, width, height, label)


# ─────────────────────────────────────────────────────────────────────────────
# STRUCTURAL HELPERS
# ─────────────────────────────────────────────────────────────────────────────

def set_white_bg(slide):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_WHITE


def accent_bar(slide):
    add_rect(slide, 0.0, 0.0, BAR_W, 5.625, C_DARK)


def footer(slide, n, total):
    add_rect(slide, 0.5, 5.08, 9.0, 0.012, C_RULE)
    add_text(slide, f"{n} / {total}",
             left=8.9, top=5.18, width=0.7, height=0.25,
             size=8.5, color=C_MUTED, align=PP_ALIGN.RIGHT)


def slide_title(slide, title):
    add_text(slide, title,
             left=0.42, top=0.15, width=9.1, height=0.68,
             size=26, bold=True, color=C_DARK)


def method_tag(slide, tag_text):
    add_rect(slide, 8.12, 0.2, 1.65, 0.28, C_DARK)
    add_text(slide, tag_text,
             left=8.12, top=0.2, width=1.65, height=0.28,
             size=9, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)


def col_divider(slide):
    add_rect(slide, 4.93, 0.95, 0.012, 3.98, C_RULE)


def col_header(slide, label, left, top, width=4.25):
    add_text(slide, label, left=left, top=top, width=width, height=0.3,
             size=10.5, bold=True, color=C_DARK)
    add_rect(slide, left, top + 0.32, width, 0.012, C_RULE)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE BUILDERS
# ─────────────────────────────────────────────────────────────────────────────

def build_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(slide)
    accent_bar(slide)

    add_text(slide,
             "MARKET RISK MODELLING  ·  FRANKFURT SCHOOL OF FINANCE & MANAGEMENT  ·  SUMMER TERM 2026",
             left=0.38, top=0.22, width=9.3, height=0.32,
             size=8, color=C_MUTED, align=PP_ALIGN.CENTER)

    add_text(slide, "VaR Calculation\n& Evolution",
             left=0.5, top=1.05, width=9.0, height=2.0,
             size=48, bold=True, color=C_DARK,
             align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide,
             "Delta-Normal  ·  Historical Simulation  ·  Monte Carlo  ·  GARCH  ·  EVT  ·  GARCH+EVT",
             left=0.5, top=3.12, width=9.0, height=0.44,
             size=12.5, color=C_MUTED, align=PP_ALIGN.CENTER)

    add_rect(slide, 3.6, 3.67, 2.8, 0.02, C_RULE)

    add_text(slide, "Dr. Sebastian Irle  ·  Group Project  ·  May 2026",
             left=0.5, top=3.85, width=9.0, height=0.38,
             size=11, color=C_MUTED, align=PP_ALIGN.CENTER)


def build_agenda(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(slide)
    accent_bar(slide)
    slide_title(slide, "Agenda")
    footer(slide, n, total)

    sections = [
        ("01", "Portfolio Overview"),
        ("02", "Methodology Overview"),
        ("03", "Delta-Normal VaR"),
        ("04", "Historical Simulation VaR"),
        ("05", "Monte Carlo VaR"),
        ("06", "GARCH-based VaR"),
        ("07", "EVT / GARCH+EVT"),
        ("08", "Assumption Validation"),
        ("09", "Backtesting & Validation"),
        ("10", "Comparative Results"),
        ("11", "Conclusion"),
    ]

    col1, col2 = sections[:6], sections[6:]
    row_h = 0.46

    def render_col(items, x_num, x_label):
        for i, (num, label) in enumerate(items):
            y = 1.05 + i * row_h
            add_rect(slide, x_num, y + 0.04, 0.36, 0.32, C_DARK)
            add_text(slide, num, left=x_num, top=y + 0.04,
                     width=0.36, height=0.32,
                     size=9, bold=True, color=C_WHITE,
                     align=PP_ALIGN.CENTER, v_anchor=MSO_ANCHOR.MIDDLE)
            add_text(slide, label, left=x_label, top=y + 0.03,
                     width=3.6, height=0.38, size=12.5, color=C_BODY)

    render_col(col1, 0.5,  0.92)
    render_col(col2, 5.15, 5.57)


def build_portfolio(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(slide)
    accent_bar(slide)
    slide_title(slide, "Portfolio Overview")
    footer(slide, n, total)
    col_divider(slide)

    add_text(slide, "Positions", left=0.5, top=1.0, width=4.25, height=0.3,
             size=11, bold=True, color=C_DARK)

    positions = [
        ("SPY",      "S&P 500 ETF",           "Equity",     "Linear",     False),
        ("IEF",      "7–10Y US Treasury ETF",  "Fixed Inc.", "Linear",     False),
        ("GLD",      "Gold ETF",               "Commodity",  "Linear",     False),
        ("EURUSD",   "EUR/USD Spot",            "FX",         "Linear",     False),
        ("IRS",      "Interest Rate Swap",      "Rates",      "Non-linear", True),
        ("Straddle", "ATM Straddle on SPY",     "Eq. Vol",    "Non-linear", True),
    ]
    cols_x = [0.5,  1.22, 2.88, 3.85]
    cols_w = [0.68, 1.62, 0.93, 1.00]
    hdrs   = ["Ticker", "Description", "Class", "Type"]
    hy     = 1.35

    add_rect(slide, 0.5, hy, 4.25, 0.3, C_DARK)
    for j, (h, x) in enumerate(zip(hdrs, cols_x)):
        add_text(slide, h, left=x + 0.06, top=hy + 0.04,
                 width=cols_w[j], height=0.22,
                 size=9, bold=True, color=C_WHITE)

    for i, (tkr, desc, cls, typ, is_nl) in enumerate(positions):
        ry = hy + 0.3 + i * 0.38
        bg = C_ROW_ALT if i % 2 == 0 else C_WHITE
        add_rect(slide, 0.5, ry, 4.25, 0.36, bg)
        c = C_DARK if is_nl else C_BODY
        for j, val in enumerate([tkr, desc, cls, typ]):
            add_text(slide, val, left=cols_x[j] + 0.06, top=ry + 0.07,
                     width=cols_w[j], height=0.24,
                     size=10, color=c, bold=is_nl)

    add_text(slide, "Sample & Data", left=5.1, top=1.0, width=4.5, height=0.3,
             size=11, bold=True, color=C_DARK)

    details = [
        ("Period",       "January 2006 – December 2024"),
        ("Horizon",      "1-day, 99% confidence level"),
        ("Risk Factors", "Asset prices, DGS10 (FRED), VIX"),
        ("Pricing",      "Full repricing — IRS (DV01), Straddle (B-S)"),
        ("Data",         "Yahoo Finance + FRED"),
        ("Est. Window",  "250-day rolling, re-estimated daily"),
    ]

    for i, (k, v) in enumerate(details):
        y = 1.35 + i * 0.48
        add_rect(slide, 5.1, y, 0.06, 0.34, C_DARK)
        add_text(slide, k + ":", left=5.24, top=y, width=1.55, height=0.34,
                 size=10.5, bold=True, color=C_DARK)
        add_text(slide, v,       left=6.82, top=y, width=2.82, height=0.34,
                 size=10.5, color=C_BODY)

    add_rect(slide, 5.1, 4.35, 4.5, 0.62, RGBColor(0xF0, 0xF0, 0xF0))
    add_text(slide,
             "The portfolio mixes linear and non-linear instruments to stress-test "
             "each VaR method — especially Delta-Normal's delta approximation.",
             left=5.2, top=4.4, width=4.3, height=0.55,
             size=9.5, italic=True, color=C_MUTED)


def build_methodology(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(slide)
    accent_bar(slide)
    slide_title(slide, "Methodology Overview")
    footer(slide, n, total)

    methods = [
        ("01", "Delta-\nNormal",    "Analytical\nNormality assumed"),
        ("02", "Historical\nSim.",  "Non-parametric\nEmpirical PnL"),
        ("03", "Monte\nCarlo",      "Simulation\nNormal + t-copula"),
        ("04", "GARCH",             "Vol clustering\nStudent-t innov."),
        ("05", "EVT",               "Tail modelling\nGPD / POT"),
        ("06", "GARCH\n+EVT",       "Combined\nDynamic tail"),
    ]

    card_w = 1.44
    gap    = 0.12
    card_h = 2.5
    card_y = 1.0

    for i, (num, name, desc) in enumerate(methods):
        x = 0.5 + i * (card_w + gap)
        add_rect(slide, x, card_y, card_w, card_h, C_CARD_BG)
        add_rect(slide, x, card_y, card_w, 0.08, C_DARK)
        add_text(slide, num,  left=x + 0.1, top=card_y + 0.15,
                 width=0.4, height=0.3, size=14, bold=True, color=C_DARK)
        add_text(slide, name, left=x + 0.1, top=card_y + 0.48,
                 width=card_w - 0.2, height=0.7,
                 size=11, bold=True, color=C_DARK)
        add_text(slide, desc, left=x + 0.1, top=card_y + 1.22,
                 width=card_w - 0.2, height=0.9,
                 size=9.5, color=C_MUTED)

    add_text(slide,
             "All six methods are applied to the same portfolio (2006–2024). "
             "Assumption validation and Kupiec / Christoffersen backtesting applied to each.",
             left=0.5, top=3.65, width=9.0, height=0.5,
             size=10.5, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    add_placeholder(slide, 0.5, 4.2, 9.0, 0.72,
                    "[ Figure: VaR time-series overview — all six methods — to be added ]")


def build_method_slide(prs, n, total,
                       title, tag,
                       bullets,
                       latex_formula,
                       formula_ref,
                       fig_name, fig_label):
    """
    Two-column VaR method slide.
    Left  : Theory bullets  +  rendered formula image  +  reference footnote
    Right : Figure or grey placeholder
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(slide)
    accent_bar(slide)
    slide_title(slide, title)
    method_tag(slide, tag)
    footer(slide, n, total)
    col_divider(slide)

    # ── Left: Theory ─────────────────────────────────────────────────────────
    col_header(slide, "Theory & Assumptions", left=0.5, top=1.0)

    add_bullets(slide, bullets,
                left=0.5, top=1.38, width=4.25, height=1.9,
                size=11.5, spacing_after=7)

    # Formula box
    add_text(slide, "Key Formula",
             left=0.5, top=3.33, width=4.25, height=0.26,
             size=9.5, bold=True, color=C_MUTED)

    add_rect(slide, 0.5, 3.61, 4.25, 0.88, C_FORMULA)

    # Render and centre the formula inside the box
    # Box centre x = 0.5 + 4.25/2 = 2.625;  vertically centred inside box
    add_formula(slide, latex_formula,
                center_x=2.625, top=3.68,
                max_w=4.1, font_size=18, dpi=200)

    add_text(slide, formula_ref,
             left=0.5, top=4.55, width=4.25, height=0.38,
             size=9, italic=True, color=C_MUTED)

    # ── Right: Results ───────────────────────────────────────────────────────
    col_header(slide, "Results", left=5.1, top=1.0, width=4.55)

    add_figure_or_ph(slide, fig_name,
                     left=5.1, top=1.38, width=4.55, height=3.58,
                     label=fig_label)


def build_assumption_validation(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(slide)
    accent_bar(slide)
    slide_title(slide, "Assumption Validation")
    footer(slide, n, total)

    add_text(slide,
             "All H\u2080 rejected across assets — confirming fat tails and volatility clustering",
             left=0.5, top=0.92, width=9.0, height=0.3,
             size=11, italic=True, color=C_MUTED)

    add_figure_or_ph(slide, "qq_plots.png",
                     left=0.5,  top=1.28, width=2.92, height=1.7,
                     label="[ QQ plots — SPY, IEF, GLD, EURUSD ]")
    add_figure_or_ph(slide, "acf_squared.png",
                     left=3.55, top=1.28, width=2.92, height=1.7,
                     label="[ ACF of squared returns — all assets ]")
    add_figure_or_ph(slide, "arch_pvalues.png",
                     left=6.6,  top=1.28, width=2.92, height=1.7,
                     label="[ Engle ARCH test p-values ]")

    add_text(slide, "Diagnostic Test Summary (p-values)",
             left=0.5, top=3.1, width=9.0, height=0.28,
             size=10.5, bold=True, color=C_DARK)

    rows = [
        ["Test",          "SPY",     "IEF",     "GLD",     "EURUSD"],
        ["Jarque-Bera",   "< 0.001", "< 0.001", "< 0.001", "< 0.001"],
        ["Ljung-Box (r)", "< 0.001", "< 0.001", "< 0.001", "< 0.001"],
        ["Engle ARCH",    "< 0.001", "< 0.001", "< 0.001", "< 0.001"],
    ]
    col_x = [0.5,  3.2,  4.72, 6.24, 7.76]
    col_w = [2.65, 1.47, 1.47, 1.47, 1.47]
    hy    = 3.43

    add_rect(slide, 0.5, hy, 9.0, 0.32, C_DARK)
    for j, (h, x) in enumerate(zip(rows[0], col_x)):
        align = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
        add_text(slide, h, left=x + 0.06, top=hy + 0.05,
                 width=col_w[j], height=0.22,
                 size=9.5, bold=True, color=C_WHITE, align=align)

    for i, row in enumerate(rows[1:]):
        ry = hy + 0.32 + i * 0.36
        bg = C_ROW_ALT if i % 2 == 0 else C_WHITE
        add_rect(slide, 0.5, ry, 9.0, 0.34, bg)
        for j, (val, x) in enumerate(zip(row, col_x)):
            align = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            add_text(slide, val, left=x + 0.06, top=ry + 0.06,
                     width=col_w[j], height=0.24,
                     size=9.5, color=C_DARK if j == 0 else C_BODY,
                     bold=(j == 0), align=align)


def build_backtesting(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(slide)
    accent_bar(slide)
    slide_title(slide, "Backtesting & Validation")
    footer(slide, n, total)
    col_divider(slide)

    col_header(slide, "Statistical Tests", left=0.5, top=1.0)

    tests = [
        "Kupiec POF Test — exception frequency matches \u03b1 = 1%?",
        "Christoffersen CC Test — additionally tests independence of exceptions",
        "Dynamic Quantile (DQ) Test — regresses hit variable on lagged predictors",
        "Traffic-light system: green / yellow / red zones (Basel III)",
        "Expected Shortfall validation via Fissler\u2013Ziegel joint elicitability",
        "Delta-Normal reference: ~100 exceptions vs. ~44 expected (fat tails confirmed)",
    ]
    add_bullets(slide, tests, left=0.5, top=1.38, width=4.25, height=3.1,
                size=11.5, spacing_after=8)

    col_header(slide, "Exception Summary", left=5.1, top=1.0, width=4.55)
    add_placeholder(slide, 5.1, 1.38, 4.55, 1.58,
                    "[ Table: Exception counts & Kupiec / CC test results per method — to be added ]")
    add_placeholder(slide, 5.1, 3.08, 4.55, 1.88,
                    "[ Figure: Exception overlay chart — all methods (2006–2024) — to be added ]")


def build_comparison(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(slide)
    accent_bar(slide)
    slide_title(slide, "Comparative Results")
    footer(slide, n, total)

    add_text(slide, "VaR Time-Series — All Six Methods",
             left=0.5, top=1.0, width=9.0, height=0.3,
             size=10.5, bold=True, color=C_DARK)
    add_rect(slide, 0.5, 1.32, 9.0, 0.012, C_RULE)

    add_placeholder(slide, 0.5, 1.38, 9.0, 2.0,
                    "[ Figure: VaR overlay — Delta-Normal, HistSim, MC, GARCH, EVT, GARCH+EVT (2006–2024) — to be added ]")

    add_text(slide, "2008–2009 Financial Crisis",
             left=0.5, top=3.5, width=4.35, height=0.26,
             size=9.5, bold=True, color=C_DARK)
    add_text(slide, "2020 COVID Crash",
             left=5.15, top=3.5, width=4.35, height=0.26,
             size=9.5, bold=True, color=C_DARK)

    add_placeholder(slide, 0.5,  3.78, 4.35, 1.18,
                    "[ Figure: GFC zoom-in — to be added ]")
    add_placeholder(slide, 5.15, 3.78, 4.35, 1.18,
                    "[ Figure: COVID crash zoom-in — to be added ]")


def build_conclusion(prs, n, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_white_bg(slide)
    accent_bar(slide)
    slide_title(slide, "Conclusion")
    footer(slide, n, total)

    bullets = [
        "Delta-Normal underestimates tail risk: ~100 exceptions vs. ~44 expected "
        "\u2014 fat tails and skewness confirmed in portfolio PnL",
        "Historical Simulation captures fat tails naturally but suffers from the "
        "\u2018ghost effect\u2019 and slow adaptation to new volatility regimes",
        "Monte Carlo with t-copula provides flexible tail dependency; "
        "normal variant still underestimates during crisis periods",
        "GARCH(1,1) with Student-t innovations correctly adapts to volatility "
        "clustering \u2014 strong improvement over static covariance",
        "EVT / GARCH+EVT delivers the most rigorous tail modelling; "
        "best Kupiec & Christoffersen test outcomes expected",
        "Non-linear instruments (IRS, Straddle) expose the failure of the delta "
        "approximation \u2014 full repricing essential for accurate VaR",
        "Method choice materially affects VaR estimates during stress periods "
        "(2008\u201309 GFC, 2020 COVID crash)",
    ]
    add_bullets(slide, bullets, left=0.5, top=1.05, width=9.0, height=3.8,
                size=12.5, spacing_after=7)

    add_text(slide,
             "Market Risk Modelling  \u00b7  Frankfurt School of Finance & Management  \u00b7  Dr. Sebastian Irle  \u00b7  May 2026",
             left=0.5, top=5.2, width=9.0, height=0.28,
             size=8.5, color=C_MUTED, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────

def build(output="var_presentation.pptx"):
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    TOTAL = 13
    print("Building slides…")

    build_title(prs)                # 1
    build_agenda(prs, 2, TOTAL)     # 2
    build_portfolio(prs, 3, TOTAL)  # 3
    build_methodology(prs, 4, TOTAL)# 4

    build_method_slide(             # 5 — Delta-Normal
        prs, 5, TOTAL,
        title="Delta-Normal VaR",
        tag="METHOD 01",
        bullets=[
            "Risk factor returns assumed jointly normal",
            "Covariance matrix \u03a3 estimated on 250-day rolling window",
            "Non-linear positions linearised via delta \u2014 gamma ignored",
            "Straddle: delta \u2248 0 at ATM \u2192 VaR \u2248 0; underestimates badly",
        ],
        latex_formula=(
            r"\mathrm{VaR}_\alpha = -\!\left(\mu_P + z_\alpha\,\sigma_P\right),"
            r"\quad \sigma_P^2 = \boldsymbol{\Delta}^\top\boldsymbol{\Sigma}\boldsymbol{\Delta}"
        ),
        formula_ref="Ref: Irle Lecture Notes, Eq. (3.1)\u2013(3.8), pp. 12\u201318",
        fig_name="delta_normal_backtest.png",
        fig_label="[ Figure: Delta-Normal VaR with P&L exceptions overlay (2006\u20132024) ]",
    )

    build_method_slide(             # 6 — Historical Simulation
        prs, 6, TOTAL,
        title="Historical Simulation VaR",
        tag="METHOD 02",
        bullets=[
            "Non-parametric: no distributional assumption required",
            "Full repricing every scenario: IRS via DV01, Straddle via Black-Scholes",
            "Captures fat tails and asymmetry naturally",
            "Ghost effect: extreme events persist exactly T = 250 days in window",
        ],
        latex_formula=(
            r"\widehat{\mathrm{VaR}}_\alpha = -\hat{Q}_\alpha"
            r"\!\left(\,\{\mathrm{PnL}_t\}_{t=1}^{T}\,\right)"
        ),
        formula_ref="Ref: Irle Lecture Notes, Eq. (4.1)\u2013(4.5), pp. 22\u201326",
        fig_name="histsim_backtest.png",
        fig_label="[ Figure: HistSim VaR with exceptions overlay \u2014 to be added ]",
    )

    build_method_slide(             # 7 — Monte Carlo
        prs, 7, TOTAL,
        title="Monte Carlo VaR",
        tag="METHOD 03",
        bullets=[
            "Simulate N = 10,000 risk-factor return scenarios per day",
            "Normal variant: Cholesky decomposition of \u03a3",
            "t-copula: marginal Student-t + Gaussian/t dependence structure",
            "Full repricing applied to each scenario",
        ],
        latex_formula=(
            r"\mathbf{r} \sim \mathcal{N}(\mathbf{0},\boldsymbol{\Sigma}),"
            r"\quad \mathrm{VaR} = -Q_\alpha\!\left(\{\tilde{P}_i\}_{i=1}^{N}\right)"
        ),
        formula_ref="Ref: Irle Lecture Notes, Eq. (5.1)\u2013(5.9), pp. 30\u201337",
        fig_name="mc_backtest.png",
        fig_label="[ Figure: MC VaR (normal & t-copula) with exceptions \u2014 to be added ]",
    )

    build_method_slide(             # 8 — GARCH
        prs, 8, TOTAL,
        title="GARCH-based VaR",
        tag="METHOD 04",
        bullets=[
            "Captures volatility clustering confirmed by Engle ARCH test",
            "Student-t innovations with estimated degrees of freedom \u03bd",
            "Estimated via rugarch (R) with ugarchroll, 250-day window",
            "VaR adapts dynamically to current volatility regime",
        ],
        latex_formula=(
            r"\sigma_t^2 = \omega + \alpha\,\varepsilon_{t-1}^2 + \beta\,\sigma_{t-1}^2,"
            r"\quad \varepsilon_t \sim t_\nu"
        ),
        formula_ref="Ref: Irle Lecture Notes, Eq. (6.1)\u2013(6.11), pp. 40\u201348",
        fig_name="garch_backtest.png",
        fig_label="[ Figure: GARCH VaR with conditional \u03c3_t time-series \u2014 to be added ]",
    )

    build_method_slide(             # 9 — EVT
        prs, 9, TOTAL,
        title="EVT / GARCH+EVT VaR",
        tag="METHOD 05",
        bullets=[
            "Peaks-over-Threshold (POT): exceedances above threshold u",
            "Threshold u selected via mean-excess plot & Hill estimator",
            "GARCH+EVT: fit GPD to GARCH standardised residuals",
            "R packages: rugarch + POT / ismev; most theoretically rigorous",
        ],
        latex_formula=(
            r"F_{\xi,\beta}(x) = 1 - \left(1 + \dfrac{\xi\,x}{\beta}\right)^{\!-1/\xi}"
        ),
        formula_ref="Ref: Irle Lecture Notes, Eq. (7.1)\u2013(7.8), pp. 52\u201360",
        fig_name="evt_backtest.png",
        fig_label="[ Figure: EVT & GARCH+EVT VaR with GPD tail fit \u2014 to be added ]",
    )

    build_assumption_validation(prs, 10, TOTAL)  # 10
    build_backtesting(prs, 11, TOTAL)             # 11
    build_comparison(prs, 12, TOTAL)              # 12
    build_conclusion(prs, 13, TOTAL)              # 13

    prs.save(output)
    print(f"\u2713  Saved  \u2192  {output}  ({TOTAL} slides)")


if __name__ == "__main__":
    build()