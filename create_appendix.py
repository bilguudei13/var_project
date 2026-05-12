"""
Generate appendix PowerPoint for MC VaR presentation.
Run from var_project root:  python create_appendix.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os
from pathlib import Path

# ── Paths ─────────────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent
FIGS = ROOT / "outputs" / "figures"
OUT  = ROOT / "outputs" / "MC_VaR_Appendix.pptx"

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x1A, 0x2B, 0x4A)   # slide header backgrounds
BLUE   = RGBColor(0x1F, 0x6F, 0xBF)   # section divider background
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF5, 0xF5, 0xF5)   # slide background
DGRAY  = RGBColor(0x44, 0x44, 0x44)   # body text
ACCENT = RGBColor(0xE8, 0x6A, 0x1F)   # orange highlights
GREEN  = RGBColor(0x1A, 0x7A, 0x3C)
RED    = RGBColor(0xC0, 0x20, 0x20)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank


# ═══════════════════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

def bg(slide, color=LGRAY):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def header_bar(slide, title, subtitle=None,
               bar_color=NAVY, left=0, top=0,
               width=None, height=Inches(1.05)):
    """Dark header bar with title (and optional subtitle)."""
    w = width or prs.slide_width
    bar = slide.shapes.add_shape(1, left, top, w, height)   # MSO_SHAPE_TYPE.RECTANGLE=1
    bar.fill.solid()
    bar.fill.fore_color.rgb = bar_color
    bar.line.fill.background()

    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.bold  = True
    run.font.size  = Pt(22)
    run.font.color.rgb = WHITE
    run.font.name  = "Calibri"

    if subtitle:
        from pptx.util import Pt as _Pt
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size  = _Pt(12)
        r2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
        r2.font.name  = "Calibri"

    # left padding
    from pptx.oxml.ns import qn
    from lxml import etree
    txBody = bar.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('lIns', str(int(Inches(0.2))))
        bodyPr.set('tIns', str(int(Inches(0.12))))

def add_textbox(slide, text, left, top, width, height,
                fontsize=Pt(11), bold=False, color=DGRAY,
                align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    # clear default paragraph
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = fontsize
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = font
    return txb

def add_lines(slide, lines, left, top, width, height,
              fontsize=Pt(11), color=DGRAY, line_spacing=None,
              bold_first=False, font="Calibri"):
    """Add multiple lines of text into one text box."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.size  = fontsize
        run.font.color.rgb = color
        run.font.name  = font
        run.font.bold  = (bold_first and i == 0)
    return txb

def add_image(slide, path, left, top, width=None, height=None):
    p = str(path)
    if not os.path.exists(p):
        return None
    try:
        if width and height:
            slide.shapes.add_picture(p, left, top, width, height)
        elif width:
            slide.shapes.add_picture(p, left, top, width=width)
        elif height:
            slide.shapes.add_picture(p, left, top, height=height)
        else:
            slide.shapes.add_picture(p, left, top)
    except Exception as e:
        print(f"  [WARN] Could not add image {p}: {e}")

def section_divider(title, subtitle=""):
    """Full-bleed blue section divider slide."""
    slide = prs.slides.add_slide(BLANK)
    bg(slide, BLUE)
    # big title
    txb = slide.shapes.add_textbox(Inches(0.6), Inches(2.5),
                                    Inches(12), Inches(1.5))
    tf = txb.text_frame
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r  = p.add_run()
    r.text = title
    r.font.size  = Pt(36)
    r.font.bold  = True
    r.font.color.rgb = WHITE
    r.font.name  = "Calibri"
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size  = Pt(16)
        r2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
        r2.font.name  = "Calibri"
    return slide

def formula_box(slide, label, formula_lines, left, top, width, height,
                label_size=Pt(11), formula_size=Pt(12)):
    """Grey box with a label and formula lines."""
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xE8, 0xEC, 0xF4)
    box.line.color.rgb = RGBColor(0xAA, 0xBB, 0xDD)

    tf = box.text_frame
    tf.word_wrap = True
    from pptx.oxml.ns import qn
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('lIns', str(int(Inches(0.12))))
        bodyPr.set('tIns', str(int(Inches(0.08))))

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = label
    r.font.size  = label_size
    r.font.bold  = True
    r.font.color.rgb = NAVY
    r.font.name  = "Calibri"

    for line in formula_lines:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = line
        r2.font.size  = formula_size
        r2.font.color.rgb = RGBColor(0x11, 0x11, 0x55)
        r2.font.name  = "Consolas"   # monospace for formulas


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE / COVER
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide, NAVY)
add_textbox(slide, "APPENDIX", Inches(0.7), Inches(1.8),
            Inches(11), Inches(0.7),
            fontsize=Pt(14), color=RGBColor(0xAA, 0xCC, 0xFF),
            bold=True, font="Calibri")
add_textbox(slide, "Monte Carlo VaR — Mathematical Reference",
            Inches(0.7), Inches(2.4), Inches(11), Inches(1.2),
            fontsize=Pt(34), bold=True, color=WHITE, font="Calibri")
add_textbox(slide, "Gaussian MC  ·  t-Copula MC  ·  GARCH-t-Copula MC",
            Inches(0.7), Inches(3.7), Inches(11), Inches(0.6),
            fontsize=Pt(16), color=RGBColor(0xCC, 0xDD, 0xFF), font="Calibri")
add_textbox(slide, "Backtesting  ·  Instrument Pricing  ·  Diagnostics",
            Inches(0.7), Inches(4.2), Inches(11), Inches(0.5),
            fontsize=Pt(16), color=RGBColor(0xCC, 0xDD, 0xFF), font="Calibri")


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION A — VAR FRAMEWORK
# ═══════════════════════════════════════════════════════════════════════════════
section_divider("A.  VaR Framework",
                "Definition  ·  Expected Shortfall  ·  Backtesting Principles")

# A1 — VaR Definition
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "A1   Value at Risk — Formal Definition",
           "The single-number loss threshold exceeded with probability (1−α)")
y = Inches(1.2)
formula_box(slide, "VaR Definition (99% confidence)",
    ["P( ΔV < −VaRₜ ) = 1 − α     where α = 0.99",
     "",
     "Equivalently:  VaRₜ = −Q_{1−α}( ΔV )",
     "",
     "  Q_{1%}(ΔV) = 1st percentile of the simulated P&L distribution",
     "  VaR is reported as a POSITIVE number (a loss amount)"],
    Inches(0.5), y, Inches(5.8), Inches(2.0))

formula_box(slide, "Expected Shortfall (ES / CVaR)",
    ["ES_{α} = −E[ ΔV  |  ΔV < −VaRₜ ]",
     "",
     "= average loss in the worst (1−α)% of scenarios",
     "",
     "ES is always ≥ VaR and is sub-additive (VaR is not)."],
    Inches(6.8), y, Inches(6.0), Inches(2.0))

formula_box(slide, "Simulation-Based Estimation (M scenarios)",
    ["Sort simulated P&Ls:  p₁ ≤ p₂ ≤ … ≤ p_M",
     "",
     "VaR  = −p_{ ⌊(1−α)M⌋ }",
     "",
     "ES   = −(1/k) Σᵢ₌₁ᵏ  pᵢ     where k = ⌊(1−α)M⌋",
     "",
     "With M=10,000 and α=99%: the 1% tail = 100 scenarios → stable estimate"],
    Inches(0.5), Inches(3.4), Inches(12.3), Inches(2.1))

add_lines(slide,
    ['Key insight: VaR answers "On a bad day (1-in-100), how much do I lose?"',
     "It does NOT tell you how bad things get beyond that threshold — that is what ES captures."],
    Inches(0.5), Inches(5.65), Inches(12.3), Inches(0.8),
    fontsize=Pt(11), color=RGBColor(0x55, 0x55, 0x55))


# A2 — Backtesting Framework
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "A2   Backtesting — Kupiec & Christoffersen Tests",
           "Statistical tests for whether observed breach frequency matches the model's confidence level")
y = Inches(1.2)
formula_box(slide, "Kupiec Proportion of Failures (POF) Test",
    ["H₀: p = 1 − α = 0.01   (true breach probability equals model assumption)",
     "",
     "Under H₀:  N ~ Binomial(T, 0.01)     E[N] = T × 0.01",
     "",
     "LR_UC = 2 [ N·ln(p̂/p₀) + (T−N)·ln((1−p̂)/(1−p₀)) ]    ~ χ²(1)",
     "",
     "  where  p̂ = N/T  (empirical breach rate),  p₀ = 0.01",
     "",
     "Reject H₀ at 95% if  LR_UC > 3.84    (i.e. N ≈ 54 for T=4019)"],
    Inches(0.5), y, Inches(12.3), Inches(2.7))

formula_box(slide, "Christoffersen Independence Test",
    ["Tests whether breaches cluster (are NOT independent) — a separate failure mode",
     "",
     "Transition counts:  n₀₀ = calm→calm,  n₀₁ = calm→breach,",
     "                    n₁₀ = breach→calm,  n₁₁ = breach→breach",
     "",
     "π₀₁ = n₀₁/(n₀₀+n₀₁)   π₁₁ = n₁₁/(n₁₀+n₁₁)   π = (n₀₁+n₁₁)/T",
     "",
     "LR_IND = 2 ln[ L(π₀₁,π₁₁) / L(π,π) ]    ~ χ²(1)",
     "",
     "Reject independence at 95% if  LR_IND > 3.84"],
    Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.7))

add_lines(slide,
    ["Our results:  Gaussian N=248, t-Copula N=25, GARCH-t-Copula N=69.",
     "Gaussian and GARCH rejected (too many breaches).  t-Copula rejected (too few — over-conservative after implementation fixes).",
     "GARCH-t-Copula is the only model that passes the Christoffersen independence test."],
    Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.45),
    fontsize=Pt(11), color=RED)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION B — PORTFOLIO
# ═══════════════════════════════════════════════════════════════════════════════
section_divider("B.  Portfolio & Data",
                "Six instruments  ·  10-year walk-forward backtest  ·  2006–2024")

# B1 — Portfolio Composition
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "B1   Portfolio Composition",
           "V₀ = $1,000,000  |  Backtest: Jan 2009 – Dec 2024  (T = 4,019 days)")

rows = [
    ("Instrument", "Type", "Weight", "Exposure", "Risk Factor"),
    ("SPY",      "US Equity ETF",    "25%", "$250,000",  "SPY log return"),
    ("GLD",      "Gold ETF",         "25%", "$250,000",  "GLD log return"),
    ("IEF",      "7-10Y Treasury",   "25%", "$250,000",  "IEF log return"),
    ("EURUSD",   "FX spot",          "25%", "$250,000",  "EURUSD log return"),
    ("IRS",      "Interest Rate Swap","—",  "$1,000,000 notional", "DGS10 yield change"),
    ("Straddle", "ATM SPY options",  "—",  "2,000 shares", "SPY price + VIX level"),
]
col_w = [Inches(1.6), Inches(2.0), Inches(1.1), Inches(2.1), Inches(5.0)]
col_x = [Inches(0.4), Inches(2.05), Inches(4.1), Inches(5.25), Inches(7.4)]
row_h = Inches(0.5)
row_y = [Inches(1.15 + i * 0.52) for i in range(len(rows))]

for r, row in enumerate(rows):
    is_header = (r == 0)
    for c, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        box = slide.shapes.add_shape(1, cx, row_y[r], cw, row_h)
        box.fill.solid()
        box.fill.fore_color.rgb = NAVY if is_header else (
            RGBColor(0xE8, 0xEC, 0xF4) if r % 2 == 0 else WHITE)
        box.line.color.rgb = RGBColor(0xAA, 0xBB, 0xDD)
        tf = box.text_frame
        from pptx.oxml.ns import qn
        bodyPr = tf._txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('lIns', str(int(Inches(0.07))))
            bodyPr.set('tIns', str(int(Inches(0.06))))
        p2 = tf.paragraphs[0]
        p2.alignment = PP_ALIGN.LEFT
        rn = p2.add_run()
        rn.text = cell
        rn.font.size = Pt(12) if not is_header else Pt(12)
        rn.font.bold = is_header
        rn.font.color.rgb = WHITE if is_header else DGRAY
        rn.font.name = "Calibri"

add_lines(slide,
    ["Linear positions: static inception shares (shares = V₀ × wⱼ / P₀ⱼ), never rebalanced.",
     "IRS: fixed payer at 3% on $1M notional, 10-year maturity.  Straddle: 30-day rolling ATM reset, 2,000 shares.",
     "Six risk factors modelled in copula: EURUSD, GLD, IEF, SPY (log returns), VIX_ret (log), DGS10_chg (absolute)."],
    Inches(0.4), Inches(4.85), Inches(12.5), Inches(0.9),
    fontsize=Pt(10), color=RGBColor(0x55, 0x55, 0x55))


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION C — GAUSSIAN MC
# ═══════════════════════════════════════════════════════════════════════════════
section_divider("C.  Gaussian Monte Carlo",
                "Multivariate normal simulation via Cholesky decomposition")

# C1 — Math
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "C1   Gaussian MC — Algorithm",
           "Rolling 750-day window  |  M = 10,000 scenarios  |  Full revaluation")

formula_box(slide, "Step 1 — Estimate parameters from rolling window W = [t−750, t)",
    ["μ  = (1/750) Σ rₜ                    (6-vector of mean returns)",
     "Σ  = (1/749) Σ (rₜ−μ)(rₜ−μ)ᵀ        (6×6 covariance matrix)",
     "Σ  += 1e-8 · I                        (regularise for Cholesky)"],
    Inches(0.5), Inches(1.15), Inches(12.3), Inches(1.35))

formula_box(slide, "Step 2 — Cholesky simulation",
    ["L  = chol(Σ)                          s.t.  L Lᵀ = Σ  (lower triangular)",
     "Z  ~ N(0, I)                          (M × 6 independent standard normals)",
     "sim = Z Lᵀ + μ                        (M × 6 correlated scenarios)",
     "",
     "Proof:  Cov(L z) = L · Cov(z) · Lᵀ = L I Lᵀ = Σ   ✓"],
    Inches(0.5), Inches(2.6), Inches(12.3), Inches(1.7))

formula_box(slide, "Step 3 — Full revaluation per scenario i",
    ["P&L_linear  = Σⱼ shares_j · price_now_j · (exp(rᵢⱼ) − 1)     j ∈ {SPY, GLD, IEF, EURUSD}",
     "P&L_IRS     = IRS(rate_now + Δrᵢ) − IRS(rate_now)",
     "P&L_strad   = [BS(S·exp(rᵢ_SPY), K, T−1/252, rₜ, σ·exp(rᵢ_VIX)) − BS(S,K,T,rₜ,σ)] × shares",
     "P&L_total   = P&L_linear + P&L_IRS + P&L_strad"],
    Inches(0.5), Inches(4.4), Inches(12.3), Inches(1.65))

add_textbox(slide, "VaR_t = −Q₁%(P&L_total)     [sign flip: VaR is a positive loss amount]",
            Inches(0.5), Inches(6.18), Inches(12.3), Inches(0.45),
            fontsize=Pt(12), bold=True, color=NAVY)


# C2 — Gaussian backtest chart
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "C2   Gaussian MC — Backtest Result",
           "N = 248 breaches  |  Expected ≈ 40  |  Kupiec REJECTED  (p ≈ 0)")
add_image(slide, FIGS / "07_mc_gaussian_var.png",
          Inches(0.4), Inches(1.15), width=Inches(12.5))
add_lines(slide,
    ["248 breaches (6.2% empirical breach rate vs 1% target).  The model systematically",
     "underestimates tail risk because fat tails and volatility clustering are both absent."],
    Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.7),
    fontsize=Pt(11), color=RED)


# C3 — Why Gaussian fails
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "C3   Why Gaussian Fails — Fat Tails",
           "Real financial returns have far more extreme events than the normal distribution predicts")

formula_box(slide, "Probability of a k-sigma event",
    ["Normal:   P(|Z| > k σ)   ↓ super-exponentially fast",
     "Student-t: P(|Z| > k σ) ~ C · k^{−ν}   (power-law tail)",
     "",
     "For k = 5 (five-sigma event):",
     "  Normal   :  P ≈ 5.7 × 10⁻⁷   →  once every ~3.5 million days",
     "  t(ν = 5) :  P ≈ 5.8 × 10⁻⁵   →  once every ~17,000 days  (206× more likely)",
     "  t(ν = 3) :  P ≈ 5.0 × 10⁻⁴   →  once every ~2,000 days   (2,850× more likely)",
     "",
     "Empirically: SPY has had several 5σ+ days within a 15-year period."],
    Inches(0.5), Inches(1.15), Inches(7.0), Inches(3.8))

add_image(slide, FIGS / "05_return_distributions.png",
          Inches(7.6), Inches(1.15), width=Inches(5.3))

formula_box(slide, "Additional failure: no volatility clustering",
    ["Gaussian model uses fixed Σ estimated over 750 days.",
     "After a crash, Σ barely changes for weeks → VaR stays low",
     "even though the market is clearly in a high-vol regime.",
     "GARCH directly addresses this by giving the model a memory."],
    Inches(0.5), Inches(5.1), Inches(7.0), Inches(1.85))


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION D — t-COPULA MC
# ═══════════════════════════════════════════════════════════════════════════════
section_divider("D.  t-Copula Monte Carlo",
                "Student-t marginals  ·  Sklar decomposition  ·  Profile MLE for ν")

# D1 — Copula theory
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "D1   Copula Theory — Sklar's Theorem",
           "Any multivariate distribution can be decomposed into marginals + a copula")

formula_box(slide, "Sklar's Theorem (1959)",
    ["For any joint CDF F(x₁,…,xₙ) with marginals F₁,…,Fₙ:",
     "",
     "  F(x₁,…,xₙ) = C( F₁(x₁), …, Fₙ(xₙ) )",
     "",
     "where C: [0,1]ⁿ → [0,1] is the COPULA — a joint CDF with uniform marginals.",
     "The copula captures the entire dependence structure, independently of the marginals.",
     "",
     "Inversion (simulation use):  if U = (U₁,…,Uₙ) ~ C  then",
     "  Xⱼ = Fⱼ⁻¹(Uⱼ)  has marginal Fⱼ  and joint dependence C."],
    Inches(0.5), Inches(1.15), Inches(12.3), Inches(2.8))

formula_box(slide, "t-Copula vs Gaussian Copula — Tail Dependence",
    ["Gaussian copula:  tail dependence λ_U = λ_L = 0  (assets become INDEPENDENT in extremes)",
     "t-Copula:         λ_L = 2·t_{ν+1}( −√((ν+1)(1−ρ)/(1+ρ)) )  > 0  for finite ν",
     "",
     "Example (ν=5, ρ=0.7):  P(SPY<1% AND GLD<1%) under t-copula ≈ 0.52%",
     "                        vs independence: 0.01%  and Gaussian copula: ≈ 0.11%",
     "",
     "This is why Gaussian copula failed in 2008 — it had zero tail dependence."],
    Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.3))


# D2 — Student-t marginals
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "D2   Student-t Marginals — Estimation",
           "Each risk factor fitted independently by maximum likelihood")

formula_box(slide, "Student-t PDF",
    ["f(x; ν, μ, σ) = Γ((ν+1)/2) / [Γ(ν/2)·√(νπ)·σ] · (1 + (x−μ)²/(ν·σ²))^{−(ν+1)/2}",
     "",
     "Parameters:  ν = degrees of freedom  (tail thickness; small ν = fatter tails)",
     "             μ = location  (≈ mean for ν > 1)",
     "             σ = scale     (≈ std dev × √((ν−2)/ν) for ν > 2)",
     "",
     "As ν → ∞:  t → Normal.   For ν ≤ 4:  variance may be infinite."],
    Inches(0.5), Inches(1.15), Inches(7.5), Inches(2.8))

formula_box(slide, "Fitted ν values (GARCH residuals, last refit)",
    ["EURUSD   :  ν ≈  7.4   (fat tails)",
     "GLD      :  ν ≈  8.3   (fat tails)",
     "IEF      :  ν ≈ 23.0   (near-Gaussian in calm periods)",
     "SPY      :  ν ≈ 11.1   (moderate tails)",
     "VIX_ret  :  ν ≈  5.0   (very fat tails + strong skew)",
     "DGS10_chg:  ν ≈ 17.2   (near-Gaussian)"],
    Inches(8.2), Inches(1.15), Inches(4.6), Inches(2.8))

formula_box(slide, "Hazen Rank Transform → Pseudo-observations",
    ["For each factor column j, compute uniform pseudo-observations:",
     "",
     "  Uᵢⱼ = rank(xᵢⱼ) / (n + 1)     (Hazen formula; avoids 0 and 1)",
     "",
     "These Uᵢⱼ ~ Uniform(0,1) marginally but retain the empirical",
     "dependence structure.  They are fed into the copula likelihood."],
    Inches(0.5), Inches(4.1), Inches(12.3), Inches(1.9))

formula_box(slide, "t-Copula Profile MLE",
    ["ν_copula = argmax_ν  Σₜ log c(U_t; R, ν)     over grid ν ∈ {2,3,…,20}",
     "R estimated by Kendall's τ inversion:  ρᵢⱼ = sin(π/2 · τ̂ᵢⱼ)"],
    Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.1))


# D3 — Simulation algorithm
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "D3   t-Copula Simulation Algorithm",
           "Cholesky + chi-squared mixing → joint fat-tailed copula draws")

formula_box(slide, "t-Copula Simulation (M scenarios, d=6 factors)",
    ["1.  L = chol(R)                     (6×6 Cholesky of correlation matrix R)",
     "2.  Z ~ N(0, I_{d×d})              (M×6 independent standard normals)",
     "3.  Y = Z Lᵀ                        (M×6 correlated normals with Cov = R)",
     "4.  W ~ χ²(ν)                        (M independent chi-squared draws)",
     "5.  T = Y · √(ν / W)               (M×6;  divide each row by √(W/ν))",
     "",
     "    Interpretation:  W is a SHARED random scaling per scenario.",
     "    When W is small (a rare event), ALL assets are scaled up together",
     "    → simultaneous extreme moves (tail dependence).",
     "",
     "6.  U = F_{t,ν}(T)                  (apply t-CDF column-by-column → uniform)"],
    Inches(0.5), Inches(1.15), Inches(12.3), Inches(3.9))

formula_box(slide, "Invert Marginals → Simulated Returns (Sklar in reverse)",
    ["7.  rᵢⱼ = F_{t, νⱼ, μⱼ, σⱼ}⁻¹(Uᵢⱼ)    for j = 1…6",
     "",
     "    Each column gets the correct marginal tail shape (factor-specific ν)",
     "    while the joint dependence is governed by the t-copula above."],
    Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.4))

formula_box(slide, "Rolling refit schedule",
    ["Refit marginals + copula every 50 days on the most recent 750 observations.",
     "Fresh M=10,000 copula draws every day (rng NOT reset between days → stochastic VaR series)."],
    Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.65))


# D4 — t-Copula backtest
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "D4   t-Copula MC — Backtest Result",
           "N = 25 breaches  |  Expected ≈ 40  |  Kupiec REJECTED  (too conservative)")
add_image(slide, FIGS / "mc_t_copula_var_backtest.png",
          Inches(0.4), Inches(1.1), width=Inches(12.5))
add_lines(slide,
    ["25 breaches — VaR is too conservative after implementation fixes (static inception shares + correct IRS/straddle pricing).",
     "Fat-tailed marginals (ν≈2–3) with no vol conditioning overstate risk once the portfolio's grown 3–4× since inception.",
     "Opposite failure mode to GARCH: too few breaches, not too many."],
    Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.7),
    fontsize=Pt(11), color=RED)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION E — GARCH-t-COPULA
# ═══════════════════════════════════════════════════════════════════════════════
section_divider("E.  GARCH-t-Copula Monte Carlo",
                "GARCH(1,1) volatility forecasting  ·  EWMA dynamic correlation  ·  Full pipeline")

# E1 — GARCH spec
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "E1   GARCH(1,1) Specification",
           "Gives the model a memory of recent volatility — the key upgrade over the t-Copula")

formula_box(slide, "GARCH(1,1) conditional variance equation",
    ["σ²_{t+1} = ω  +  α · ε²_t  +  β · σ²_t",
     "",
     "  ω  (omega)  :  long-run variance floor.  Ensures σ²>0 at all times.",
     "                 Unconditional variance:  σ²_∞ = ω / (1 − α − β)",
     "",
     "  α  (alpha)  :  ARCH coefficient.  Size of yesterday's shock → today's vol.",
     "                 Large α = fast reaction to new information.",
     "",
     "  β  (beta)   :  GARCH coefficient.  Persistence of the vol regime.",
     "                 Large β = vol decays slowly (clustering).",
     "",
     "  α + β < 1   :  stationarity condition.  Typical empirical values: 0.97–0.99.",
     "",
     "  ν  (nu)     :  degrees of freedom of t-distributed innovations.",
     "                 Controls tail thickness of the residual distribution."],
    Inches(0.5), Inches(1.15), Inches(7.8), Inches(4.9))

formula_box(slide, "GARCH returns model",
    ["rₜ = μ + σₜ · zₜ     where  zₜ ~ t(ν) standardised",
     "",
     "Standardised t: Var(zₜ) = 1  (not ν/(ν−2))",
     "Requires std_scale = √((ν−2)/ν) correction when using scipy."],
    Inches(8.6), Inches(1.15), Inches(4.2), Inches(1.9))

formula_box(slide, "One-step-ahead forecast (used in simulation)",
    ["Given today's state (x_t, σ_t):",
     "",
     "  σ²_{t+1} = ω  +  α · x²_t  +  β · σ²_t",
     "",
     "  r_sim = μ + σ_{t+1} · z     z drawn from t-copula"],
    Inches(8.6), Inches(3.25), Inches(4.2), Inches(2.0))

formula_box(slide, "Numerical implementation note",
    ["Returns are scaled ×100 before GARCH fitting (GARCH_SCALE=100).",
     "Optimizer works in percentage-point units (e.g. 0.8% → 0.8).",
     "After fitting: ω is divided back by 100² to restore decimal units."],
    Inches(8.6), Inches(5.4), Inches(4.2), Inches(1.6))


# E2 — EWMA correlation
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "E2   EWMA Dynamic Correlation",
           "Correlation updated daily using exponentially weighted standardised residuals  (λ = 0.94)")

formula_box(slide, "EWMA Correlation Update (JP Morgan RiskMetrics)",
    ["Each day, after observing the new GARCH residual  z_t = x_t / σ_t:",
     "",
     "  Q_{t+1}  =  λ · Q_t  +  (1−λ) · z_t · z_tᵀ",
     "",
     "  R_{t+1}  =  diag(Q_{t+1})^{−1/2} · Q_{t+1} · diag(Q_{t+1})^{−1/2}",
     "",
     "  where:  Q_t  = quasi-correlation matrix (not unit diagonal)",
     "          R_t  = correlation matrix  (unit diagonal enforced by normalisation)",
     "          λ=0.94  = decay factor  (half-life ≈ 12 trading days)",
     "",
     "Why standardise by σ_t first?  We want pure correlation (not covariance).",
     "Using raw returns would mix correlation changes with volatility changes."],
    Inches(0.5), Inches(1.15), Inches(12.3), Inches(3.5))

formula_box(slide, "Warm-up loop (prevents discontinuity at each GARCH refit)",
    ["At each refit: GARCH params (ω, α, β) change → σ_t series changes → z_t series changes.",
     "To prevent a jump in Q when params change, the EWMA is replayed from scratch",
     "by running the update equation through the full 750-day window of new residuals.",
     "This ensures R_dynamic is continuous even after a refit."],
    Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.7))

add_lines(slide,
    ["Empirical correlations at last refit (Dec 2024):  ρ(SPY,IEF)=+0.12  (was −0.40 pre-2022)",
     "ρ(SPY,VIX)=−0.75  |  ρ(IEF,DGS10)=−0.96  |  ρ(GLD,IEF)=+0.43"],
    Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.75),
    fontsize=Pt(11), color=DGRAY)


# E3 — Full pipeline
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "E3   GARCH-t-Copula — Complete 8-Step Pipeline",
           "Every step listed with inputs/outputs")

steps = [
    ("Step 1", "Fit GARCH(1,1)-t per factor",
     "Input: 750-day rolling window returns. Output: (ω,α,β,μ,ν) per factor."),
    ("Step 2", "Compute GARCH residuals",
     "zₜ = (rₜ − μ) / σₜ   →  standardised residuals, roughly t(ν)-distributed."),
    ("Step 3", "Rank pseudo-observations",
     "Uᵢⱼ = rank(zᵢⱼ)/(n+1).  Fed to copula likelihood for R and ν_copula estimation."),
    ("Step 3b", "Initialise EWMA",
     "Q₀ = sample cov of GARCH residuals.  Warm-up replay over 750 days."),
    ("Step 4", "Update R_dynamic daily",
     "Q_{t+1} = λQ_t + (1−λ)z_t z_tᵀ  →  R_{t+1} (normalise diagonal to 1)."),
    ("Step 5", "Simulate M t-copula draws",
     "U ~ t-Copula(R_dynamic, ν_copula).  M=10,000 uniform scenarios."),
    ("Step 6", "Invert & reconstruct returns",
     "z_sim = t_ppf(U, ν) × std_scale.   r_sim = μ + σ_forecast × z_sim."),
    ("Step 7", "Full revaluation",
     "P&L = linear (static shares) + IRS (annuity) + straddle (Black-Scholes)."),
    ("Step 8", "Extract VaR",
     "VaR_t = −Q₁%(P&L_sim).   Store and advance one day."),
]
x1, x2, x3 = Inches(0.35), Inches(1.1), Inches(3.0)
y0 = Inches(1.15)
dy = Inches(0.64)
for i, (tag, title, detail) in enumerate(steps):
    y = y0 + i * dy
    # tag pill
    pill = slide.shapes.add_shape(1, x1, y, Inches(0.6), Inches(0.48))
    pill.fill.solid()
    pill.fill.fore_color.rgb = NAVY
    pill.line.fill.background()
    tf = pill.text_frame
    from pptx.oxml.ns import qn
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        bodyPr.set('tIns', str(int(Inches(0.06))))
    p3 = tf.paragraphs[0]; p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = tag; r3.font.size = Pt(8)
    r3.font.bold = True; r3.font.color.rgb = WHITE; r3.font.name = "Calibri"
    # title
    txb2 = slide.shapes.add_textbox(x2, y, Inches(1.75), Inches(0.48))
    tf2 = txb2.text_frame; p4 = tf2.paragraphs[0]
    r4 = p4.add_run(); r4.text = title; r4.font.size = Pt(10)
    r4.font.bold = True; r4.font.color.rgb = NAVY; r4.font.name = "Calibri"
    # detail
    txb3 = slide.shapes.add_textbox(x3, y, Inches(9.9), Inches(0.48))
    tf3 = txb3.text_frame; p5 = tf3.paragraphs[0]
    r5 = p5.add_run(); r5.text = detail; r5.font.size = Pt(10)
    r5.font.color.rgb = DGRAY; r5.font.name = "Calibri"


# E4 — std_scale correction
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "E4   std_scale Correction — Unit-Variance t Draws",
           "Critical detail: scipy's t.ppf has variance ν/(ν−2), not 1")

formula_box(slide, "The problem",
    ["scipy.stats.t.ppf(u, df=ν) returns X ~ t(ν) with Var(X) = ν/(ν−2).",
     "",
     "But GARCH is calibrated under the convention Var(z_t) = 1.",
     "If we use raw ppf draws: r_sim = σ_forecast × ppf(u)",
     "  then  Var(r_sim) = σ²_forecast × ν/(ν−2)  ≠  σ²_forecast",
     "",
     "For ν=8:  ν/(ν−2) = 1.333  →  volatility overstated by 15%."],
    Inches(0.5), Inches(1.15), Inches(12.3), Inches(2.5))

formula_box(slide, "The fix",
    ["std_scale_j = √( (νⱼ − 2) / νⱼ )",
     "",
     "z_std = t.ppf(U, df=ν) × std_scale",
     "",
     "Verification:  Var(z_std) = [ν/(ν−2)] × [(ν−2)/ν] = 1   ✓",
     "",
     "r_sim = μ + σ_forecast × z_std    →  Var(r_sim) = σ²_forecast   ✓",
     "",
     "Applied element-wise across all 6 factors with their own νⱼ."],
    Inches(0.5), Inches(3.8), Inches(12.3), Inches(2.5))

add_lines(slide,
    ["This correction is NOT applied in the t-Copula model (mc_t_copula.py uses direct scipy t.ppf with loc/scale parameters).",
     "Only the GARCH model needs it because it constructs z_std explicitly and multiplies by σ_forecast separately."],
    Inches(0.5), Inches(6.45), Inches(12.3), Inches(0.8),
    fontsize=Pt(10), color=RGBColor(0x55, 0x55, 0x55))


# E5 — GARCH backtest
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "E5   GARCH-t-Copula MC — Backtest Result",
           "N = 69 breaches  |  Expected ≈ 40  |  Kupiec REJECTED  |  Independence: RETAINED")
add_image(slide, FIGS / "11_exceptions_timeline_MC-GARCH-t-Copula.png",
          Inches(0.4), Inches(1.1), width=Inches(8.5))
add_image(slide, FIGS / "13_transition_matrix_MC-GARCH-t-Copula.png",
          Inches(9.1), Inches(1.1), width=Inches(3.9))
add_lines(slide,
    ["69 breaches — a 72% reduction from Gaussian (248). All 10 worst days are in 2022 (rate-hike regime).",
     "Independence test RETAINED: breaches are not clustered, meaning the model tracks vol regimes correctly.",
     "Kupiec still rejected because 2022 was a structural regime break the rolling-window model could not anticipate."],
    Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.8),
    fontsize=Pt(11), color=DGRAY)


# E6 — GARCH stability
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "E6   GARCH Parameter Stability Over Time",
           "How ω, α, β, ν evolve across the 15-year backtest window")
add_image(slide, FIGS / "validate_mc_garch_t_copula_07_garch_stability.png",
          Inches(0.4), Inches(1.1), height=Inches(5.8))

# E7 — Copula tail validation
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "E7   Copula Tail Dependence — Model vs Empirical",
           "Simulated joint tail frequencies vs actual co-crash rates in historical data")
add_image(slide, FIGS / "validate_mc_garch_t_copula_05_copula_tail.png",
          Inches(0.4), Inches(1.1), height=Inches(5.8))

# E8 — GARCH Q-Q
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "E8   GARCH Residual Q-Q Plots",
           "Standardised residuals should lie on the t(ν) theoretical line if GARCH is correctly specified")
add_image(slide, FIGS / "validate_mc_garch_t_copula_02_garch_qq.png",
          Inches(0.4), Inches(1.1), height=Inches(5.8))

# E9 — Component P&L
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "E9   Component P&L Breakdown — Simulated vs Actual",
           "Linear, IRS, and straddle contributions to total P&L across the backtest period")
add_image(slide, FIGS / "validate_mc_garch_t_copula_09_component_pnl.png",
          Inches(0.4), Inches(1.1), height=Inches(5.8))

# E10 — Stress scenarios
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "E10   Top 10 Loss Days — 2022 Regime Break",
           "All 10 worst portfolio days occurred in 2022 — the Fed rate-hiking cycle")

rows2 = [
    ("Rank", "Date", "Total P&L", "Linear P&L", "IRS P&L", "Straddle P&L"),
    ("1",  "2022-09-13", "−$69,012",  "−$87,874", "+$3,637",  "+$15,225"),
    ("2",  "2022-06-13", "−$66,402",  "−$103,156","+$20,135", "+$16,618"),
    ("3",  "2022-05-09", "−$60,256",  "−$66,143", "−$5,123",  "+$11,010"),
    ("4",  "2022-08-26", "−$56,744",  "−$70,179", "+$729",    "+$12,706"),
    ("5",  "2022-10-07", "−$55,510",  "−$63,349", "+$4,362",  "+$3,477"),
    ("6",  "2022-12-15", "−$55,502",  "−$57,049", "−$3,655",  "+$5,202"),
    ("7",  "2022-05-18", "−$51,364",  "−$60,645", "−$6,593",  "+$15,873"),
    ("8",  "2022-04-29", "−$50,765",  "−$64,664", "+$2,911",  "+$10,987"),
    ("9",  "2022-10-14", "−$50,323",  "−$53,286", "+$2,184",  "+$779"),
    ("10", "2022-11-02", "−$46,773",  "−$50,319", "+$2,184",  "+$1,362"),
]
col_w2 = [Inches(0.55), Inches(1.3), Inches(1.7), Inches(1.8), Inches(1.5), Inches(1.8)]
col_x2 = [Inches(0.3), Inches(0.9), Inches(2.25), Inches(4.0), Inches(5.85), Inches(7.4)]
row_h2 = Inches(0.485)
for r, row in enumerate(rows2):
    y = Inches(1.15) + r * row_h2
    is_hdr = (r == 0)
    for c, (cell, cx, cw) in enumerate(zip(row, col_x2, col_w2)):
        box = slide.shapes.add_shape(1, cx, y, cw, row_h2)
        box.fill.solid()
        box.fill.fore_color.rgb = (NAVY if is_hdr else
                                   (RGBColor(0xFF, 0xEE, 0xEE) if r > 0 else WHITE))
        box.line.color.rgb = RGBColor(0xAA, 0xBB, 0xDD)
        tf = box.text_frame
        from pptx.oxml.ns import qn
        bodyPr = tf._txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('lIns', str(int(Inches(0.07))))
            bodyPr.set('tIns', str(int(Inches(0.08))))
        p2 = tf.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
        rn = p2.add_run(); rn.text = cell
        rn.font.size = Pt(11); rn.font.bold = is_hdr
        rn.font.color.rgb = WHITE if is_hdr else (RED if "−" in cell and c >= 2 else GREEN if "+" in cell else DGRAY)
        rn.font.name = "Calibri"

add_lines(slide,
    ["Key insight: linear equity positions drove the losses (stocks + bonds falling together).",
     "The straddle partially hedged through vol spikes but could not offset the scale of equity drawdowns.",
     "The IRS was a mixed hedge: rising rates hurt the fixed-payer position on some days but not others."],
    Inches(0.3), Inches(6.55), Inches(9.1), Inches(0.85),
    fontsize=Pt(10), color=DGRAY)

add_image(slide, FIGS / "validate_mc_garch_t_copula_10_stress_scenarios.png",
          Inches(9.4), Inches(1.15), width=Inches(3.5))


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION F — INSTRUMENT PRICING
# ═══════════════════════════════════════════════════════════════════════════════
section_divider("F.  Instrument Pricing",
                "Black-Scholes straddle  ·  IRS discrete annuity  ·  Consistency across models")

# F1 — Black-Scholes
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "F1   Black-Scholes ATM Straddle Pricing",
           "Full revaluation per scenario: straddle = call + put at the same strike")

formula_box(slide, "Black-Scholes Call & Put Formulas",
    ["d₁ = [ ln(S/K) + (r + σ²/2)·T ] / (σ√T)",
     "d₂ = d₁ − σ√T",
     "",
     "Call  =  S·N(d₁)  −  K·e^{−rT}·N(d₂)",
     "Put   =  K·e^{−rT}·N(−d₂)  −  S·N(−d₁)",
     "",
     "Straddle  =  Call + Put"],
    Inches(0.5), Inches(1.15), Inches(6.0), Inches(2.8))

formula_box(slide, "Variables in simulation",
    ["S      = S_now · exp(r_SPY)       (simulated next-day SPY price)",
     "σ      = σ_now · exp(r_VIX)       (simulated next-day implied vol)",
     "r      = r_now + Δr_DGS10         (simulated next-day yield)",
     "K      = strike (ATM at last 30-day reset)",
     "T      = T_now − 1/252             (one trading day of time decay)",
     "",
     "P&L_straddle = [Straddle(S_sim,K,T_next,r_sim,σ_sim)",
     "                − Straddle(S_now,K,T_now,r_now,σ_now)] × shares"],
    Inches(6.8), Inches(1.15), Inches(6.0), Inches(2.8))

formula_box(slide, "Straddle Greeks (why it profits from crises)",
    ["Vega  = ∂(Straddle)/∂σ = S·√T·φ(d₁) > 0",
     "   → gains value when implied volatility spikes (VIX surge in a crash)",
     "",
     "Delta = ∂(Straddle)/∂S = N(d₁) − N(−d₁)",
     "   → near zero ATM; small directional exposure",
     "",
     "Theta = time decay < 0  (costs money every day without a large move)",
     "",
     "Net effect in a crisis:  large Vega gains >> Theta cost + small Delta loss",
     "Net effect in calm markets:  Theta bleeds value slowly each day"],
    Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.4))


# F2 — IRS
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "F2   Interest Rate Swap — Discrete Annuity Pricing",
           "Fixed payer: receives floating (LIBOR/OIS), pays fixed 3%  |  Notional $1M, 10-year")

formula_box(slide, "Correct Formula: Discrete Annuity",
    ["V_IRS(r_float) = Notional × (r_float − r_fixed) × Σₖ₌₁ᴺ  1/(1 + r_float)ᵏ",
     "",
     "For annual payments, N = maturity (years):",
     "",
     "  Annuity factor A(r, N) = Σₖ₌₁ᴺ  1/(1+r)ᵏ  =  [1 − (1+r)^{−N}] / r",
     "",
     "  V_IRS = Notional × (r_float − r_fixed) × A(r_float, N)",
     "",
     "Positive when r_float > r_fixed (rising rates benefit fixed payer)."],
    Inches(0.5), Inches(1.15), Inches(12.3), Inches(2.8))

formula_box(slide, "Simplified Formula (INCORRECT — used in earlier models)",
    ["V_IRS_simplified = Notional × (r_float − r_fixed) × N / (1 + r_float)",
     "",
     "This applies a single discount factor to the entire undiscounted stream.",
     "At r=4%, N=10:  A_correct ≈ 8.11   vs   N/(1+r) ≈ 9.62   →  ~18% overestimate.",
     "",
     "Fixed in all three MC models: now all use portfolio_pricing.price_irs (discrete annuity)."],
    Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.1))

add_lines(slide,
    ["IRS P&L per scenario:  P&L_IRS = V_IRS(rate_sim) − V_IRS(rate_now)",
     "where rate_sim = rate_now + Δr_DGS10_sim  (absolute yield change from simulation)"],
    Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.8),
    fontsize=Pt(11), color=NAVY, bold_first=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION G — MODEL COMPARISON & DIAGNOSIS
# ═══════════════════════════════════════════════════════════════════════════════
section_divider("G.  Model Comparison & Regime Break",
                "What improved and what didn't — and why the fixes backfired")

# G1 — Summary table
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "G1   Model Comparison Summary",
           "Three models, same portfolio, same 4,019-day backtest, same α=99%")

cols  = ["Model", "Breaches N", "Expected", "Kupiec", "Independence", "Key feature added"]
widths = [Inches(2.2), Inches(1.3), Inches(1.1), Inches(1.3), Inches(1.5), Inches(5.4)]
xs    = [Inches(0.3), Inches(2.55), Inches(3.9), Inches(5.05), Inches(6.4), Inches(7.95)]
data  = [
    cols,
    ["Gaussian MC",      "248",  "≈ 40", "REJECTED", "REJECTED", "Cholesky of rolling Σ, normal draws"],
    ["t-Copula MC",      " 25",  "≈ 40", "REJECTED", "REJECTED", "+ fat-tailed marginals + tail dependence copula"],
    ["GARCH-t-Copula",   " 69",  "≈ 40", "REJECTED", "RETAINED", "+ GARCH vol memory + EWMA dynamic correlation"],
    ["Critical value",  "≈ 54",  "≈ 40", "3.84 (χ²)", "3.84 (χ²)", "Kupiec 95% threshold for T=4,019"],
]
rh = Inches(0.56)
for r, row in enumerate(data):
    y = Inches(1.15) + r * rh
    is_hdr = (r == 0)
    is_crit = (r == 4)
    for c, (cell, cx, cw) in enumerate(zip(row, xs, widths)):
        fill_c = (NAVY if is_hdr else
                  RGBColor(0xDD, 0xEE, 0xFF) if is_crit else
                  (RGBColor(0xFF, 0xEE, 0xEE) if r in (1,2,3) and c == 3 else
                   RGBColor(0xDD, 0xFF, 0xDD) if r == 3 and c == 4 else
                   (LGRAY if r % 2 == 0 else WHITE)))
        box = slide.shapes.add_shape(1, cx, y, cw, rh)
        box.fill.solid(); box.fill.fore_color.rgb = fill_c
        box.line.color.rgb = RGBColor(0xAA, 0xBB, 0xDD)
        tf = box.text_frame
        from pptx.oxml.ns import qn
        bodyPr = tf._txBody.find(qn('a:bodyPr'))
        if bodyPr is not None:
            bodyPr.set('lIns', str(int(Inches(0.07))))
            bodyPr.set('tIns', str(int(Inches(0.1))))
        p2 = tf.paragraphs[0]; p2.alignment = PP_ALIGN.LEFT
        rn = p2.add_run(); rn.text = cell
        rn.font.size = Pt(11.5)
        rn.font.bold = is_hdr or is_crit
        rn.font.color.rgb = (WHITE if is_hdr else
                             RED if ("REJECTED" in cell) else
                             GREEN if "RETAINED" in cell else DGRAY)
        rn.font.name = "Calibri"

add_lines(slide,
    ["GARCH is the only model with a passing independence test — it correctly tracks when the market is calm vs volatile.",
     "The remaining excess breaches (69 vs ~54 threshold) are concentrated in 2022 and cannot be eliminated by parameter tuning."],
    Inches(0.3), Inches(4.1), Inches(12.5), Inches(0.8),
    fontsize=Pt(11), color=DGRAY)

formula_box(slide, "Why the attempted fixes (ν-cap + skewed-t for VIX) made things WORSE",
    ["GARCH residuals → pseudo-observations → copula fit → joint simulation → straddle P&L",
     "Changing VIX distribution → reshuffles pseudo-obs → changes copula → more extreme upward VIX scenarios",
     "More extreme VIX → long straddle makes more money in simulation → VaR goes DOWN → more real breaches.",
     "The model had compensating effects around its own misspecifications.  N went 69 → 72.  Changes reverted."],
    Inches(0.3), Inches(5.05), Inches(12.5), Inches(1.7))


# G2 — 2022 Regime
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "G2   2022 Regime Break — Stocks and Bonds Fall Together",
           "The root cause of residual Kupiec failure: a structural correlation regime shift")

formula_box(slide, "Historical SPY–IEF correlation by regime",
    ["2009–2014  (post-GFC, ZIRP)  :  ρ(SPY, IEF) ≈ −0.40  to  −0.45   (bonds hedge stocks)",
     "2015–2019  (normalisation)   :  ρ(SPY, IEF) ≈ −0.35  to  −0.40   (still negative)",
     "2020–2021  (COVID / ZIRP 2)  :  ρ(SPY, IEF) ≈ −0.45               (flight-to-safety intact)",
     "2022–2023  (rate-hike cycle) :  ρ(SPY, IEF) ≈ +0.10  to  +0.30   (BREAK: positive correlation!)",
     "2024       (late cycle)      :  ρ(SPY, IEF) ≈ +0.12               (still elevated)",
     "",
     "The Fed raised rates from 0.25% to 5.50% in 18 months — fastest hiking cycle since the 1980s.",
     "Rising rates → bond prices fall (IEF down) AND equity valuations fall (SPY down) simultaneously.",
     "The model, trained on 750 days of mostly-negative SPY-IEF correlation, treated IEF as a hedge."],
    Inches(0.5), Inches(1.15), Inches(12.3), Inches(3.3))

formula_box(slide, "Why no rolling-window model could have anticipated this",
    ["The regime break in correlation happened suddenly in early 2022.",
     "A 750-day rolling window in late 2021 still showed ρ ≈ −0.40.",
     "As 2022 data flowed in, the EWMA gradually updated R_dynamic — but by then",
     "the damage to the backtest score had already accumulated.",
     "",
     "Resolution:  Stressed VaR overlay (Basel 2.5) — take the maximum of:",
     "  (a) rolling VaR (current regime),  (b) VaR from a fixed 2022 stress window.",
     "This does not eliminate the break but ensures the model is never less than the stress-period VaR."],
    Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.3))


# ═══════════════════════════════════════════════════════════════════════════════
# SECTION H — EXTRA DIAGNOSTICS
# ═══════════════════════════════════════════════════════════════════════════════
section_divider("H.  Additional Diagnostic Charts",
                "GARCH ACF  ·  Sign bias  ·  P&L distribution  ·  Correlation matrix")

for fig, title, subtitle in [
    ("validate_mc_garch_t_copula_03_garch_acf_sq.png",
     "H1   GARCH ACF — Squared Residuals",
     "No remaining autocorrelation in squared residuals confirms GARCH correctly captures volatility clustering"),
    ("validate_mc_garch_t_copula_08_sign_bias.png",
     "H2   Sign Bias Test",
     "Tests whether positive vs negative shocks have asymmetric impact on volatility"),
    ("validate_mc_garch_t_copula_01_pnl_distribution.png",
     "H3   Simulated P&L Distribution",
     "Left tail of the 10,000-scenario P&L distribution — the VaR is the 1st percentile"),
    ("08_correlation_matrix.png",
     "H4   Empirical Correlation Matrix",
     "Rolling-average pairwise correlations across the 6 risk factors"),
]:
    slide = prs.slides.add_slide(BLANK)
    bg(slide)
    header_bar(slide, title, subtitle)
    add_image(slide, FIGS / fig, Inches(0.4), Inches(1.1), height=Inches(5.9))


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"\nSaved -> {OUT}")
print(f"Slides: {len(prs.slides)}")
